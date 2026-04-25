from pptx import Presentation
from pptx.util import Inches, Emu

ppt = Presentation("项目汇报.pptx")

slide_w = ppt.slide_width
slide_h = ppt.slide_height
print(f"Slide size: {Emu(slide_w).inches:.3f}\" x {Emu(slide_h).inches:.3f}\"")

for i, slide in enumerate(ppt.slides, 1):
    print(f"\n=== Slide {i} ===")
    for shape in slide.shapes:
        left = Emu(shape.left).inches
        top = Emu(shape.top).inches
        width = Emu(shape.width).inches
        height = Emu(shape.height).inches
        right = left + width
        bottom = top + height
        text_preview = ""
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()[:40]
            text_preview = f" | text: '{text}'"
        print(f"  Shape: {shape.shape_type.name} | ({left:.2f}, {top:.2f}) {width:.2f}x{height:.2f}{text_preview}")
        if left < 0 or top < 0 or right > Emu(slide_w).inches or bottom > Emu(slide_h).inches:
            print(f"    ⚠️ OUT OF BOUNDS")
        if width < 0.3 and shape.has_text_frame and len(shape.text_frame.text) > 10:
            print(f"    ⚠️ VERY NARROW TEXT BOX")
